"""
PPTX parser.

Handles: .pptx, .ppt
Uses python-pptx to iterate slides and shapes, producing
slide-aware DocumentBlocks with page_index = slide index.
"""

from __future__ import annotations

import io
from typing import Protocol, cast

import pptx
import structlog
from django.core.files.storage import default_storage
from pptx.presentation import Presentation as PresentationDocument
from pptx.shapes.base import BaseShape
from pptx.table import Table as PptxTable

from apps.docusafe.services.parsers.document_block import DocumentBlock

logger = structlog.getLogger("default")


class TableShape(Protocol):
    table: PptxTable


def parse_pptx(file_path: str) -> list[DocumentBlock]:
    """
    Parse a PowerPoint file into slide-aware DocumentBlocks.

    Each slide is treated as a separate page (page_index = slide index).
    Title shapes become heading blocks; content shapes become paragraphs;
    table shapes are serialised as markdown tables.
    """
    prs = _load_pptx_from_s3(file_path)
    blocks: list[DocumentBlock] = []

    for slide_idx, slide in enumerate(prs.slides):
        current_section: str | None = None
        table_counter = 0

        for shape in slide.shapes:
            if shape.has_table:
                table_text = _shape_table_to_markdown(shape)
                if table_text.strip():
                    blocks.append(
                        DocumentBlock(
                            text=table_text,
                            block_type="table",
                            page_index=slide_idx,
                            section=current_section,
                            metadata={
                                "slide_number": slide_idx + 1,
                                "table_index": table_counter,
                            },
                        )
                    )
                    table_counter += 1

            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                if _is_title_shape(shape):
                    current_section = text
                    blocks.append(
                        DocumentBlock(
                            text=text,
                            block_type="heading",
                            page_index=slide_idx,
                            section=current_section,
                            metadata={
                                "slide_number": slide_idx + 1,
                                "heading_level": 1,
                            },
                        )
                    )
                else:
                    blocks.append(
                        DocumentBlock(
                            text=text,
                            block_type="paragraph",
                            page_index=slide_idx,
                            section=current_section,
                            metadata={"slide_number": slide_idx + 1},
                        )
                    )

    return blocks


def _is_title_shape(shape: BaseShape) -> bool:
    """Check if a shape is a title placeholder."""
    if shape.is_placeholder:
        # Placeholder types 0 (TITLE) and 1 (CENTER_TITLE) indicate title shapes
        ph_idx = shape.placeholder_format.idx
        return ph_idx in {0, 1}
    return False


def _shape_table_to_markdown(shape: BaseShape) -> str:
    """Convert a PPTX table shape to pipe-delimited markdown."""
    table = cast(TableShape, shape).table
    rows = []

    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if len(rows) >= 1:
        col_count = len(table.rows[0].cells)
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        rows.insert(1, separator)

    return "\n".join(rows)


def _load_pptx_from_s3(file_path: str) -> PresentationDocument:
    """Download a PPTX file from S3 and return a python-pptx Presentation."""
    try:
        with default_storage.open(file_path, "rb") as file_obj:
            content = file_obj.read()
        return pptx.Presentation(io.BytesIO(content))
    except Exception:
        logger.exception("Failed to load PPTX from S3", file_path=file_path)
        raise
